"""
office_extractors.py
====================
Structure-preserving Markdown extraction from Word, PowerPoint, and Excel files,
optimized for section-aware chunking and LLM/RAG consumption.

Classes
-------
WordExtractor(file_path, tags_dir)   -> .doc / .docx   -> page_text[page_number] = md
PPTExtractor(file_path, tags_dir)    -> .ppt / .pptx   -> page_text[slide_number] = md
ExcelExtractor(file_path, tags_dir)  -> .csv/.xls/.xlsx -> page_text[sheet_number] = md

Conventions
-----------
* Images / charts / visual objects are saved into `tags_dir` and replaced in-place by
  `## Figure-{page|slide|sheet}-{n}` placeholders.
* Tables are rendered as GitHub Markdown tables, preceded by `## Table-{page|slide|sheet}-{n}`.
* Headings, lists, checklists, code blocks, hyperlinks, footnotes, captions, bold/italic
  and reading order are preserved; noise (empty runs, duplicated headers/footers,
  redundant whitespace) is removed.

Dependencies: python-docx, python-pptx, openpyxl. Legacy formats (.doc/.ppt/.xls)
are converted in a temp dir via LibreOffice (`soffice`), which must be on PATH.
"""

from __future__ import annotations

import csv
import os
import re
import shutil
import subprocess
import tempfile
from typing import Dict, List, Optional

# --------------------------------------------------------------------------- #
# Shared helpers
# --------------------------------------------------------------------------- #

_WS_RE = re.compile(r"[ \t]+")
_NL_RE = re.compile(r"\n{3,}")
_PIPE_RE = re.compile(r"\|")


def _clean_cell(text: str) -> str:
    """Sanitize a value for a Markdown table cell."""
    text = _WS_RE.sub(" ", (text or "").replace("\n", " ").strip())
    return _PIPE_RE.sub(r"\\|", text)


def _md_table(rows: List[List[str]]) -> str:
    """Render a list of rows as a GitHub Markdown table (first row = header)."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [[_clean_cell(c) for c in r] + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join(out)


def _normalize_md(text: str) -> str:
    """Collapse extra blank lines / trailing spaces into clean Markdown."""
    lines = [ln.rstrip() for ln in text.split("\n")]
    return _NL_RE.sub("\n\n", "\n".join(lines)).strip()


def _save_image(blob: bytes, tags_dir: str, stem: str, ext: str) -> str:
    """Persist an image blob under tags_dir; returns the saved path."""
    os.makedirs(tags_dir, exist_ok=True)
    ext = (ext or "png").lstrip(".").lower() or "png"
    path = os.path.join(tags_dir, f"{stem}.{ext}")
    with open(path, "wb") as fh:
        fh.write(blob)
    return path


def _convert_with_soffice(file_path: str, target_ext: str) -> str:
    """Convert a legacy Office file to its OOXML equivalent via LibreOffice."""
    tmp = tempfile.mkdtemp(prefix="oxt_")
    subprocess.run(
        ["soffice", "--headless", "--convert-to", target_ext,
         "--outdir", tmp, file_path],
        check=True, capture_output=True, timeout=180,
        env={**os.environ, "HOME": tmp},
    )
    base = os.path.splitext(os.path.basename(file_path))[0]
    out = os.path.join(tmp, f"{base}.{target_ext}")
    if not os.path.exists(out):
        raise RuntimeError(f"LibreOffice failed to convert {file_path!r}")
    return out


# --------------------------------------------------------------------------- #
# Word
# --------------------------------------------------------------------------- #

class WordExtractor:
    """Extract .doc/.docx to per-page Markdown with figure/table tags."""

    _NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    def __init__(self, file_path: str, tags_dir: str):
        self.file_path = file_path
        self.tags_dir = tags_dir
        self.page_text: Dict[int, str] = {}

    # -- public ------------------------------------------------------------ #

    def extract(self) -> Dict[int, str]:
        import docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        path = self.file_path
        if path.lower().endswith(".doc"):
            path = _convert_with_soffice(path, "docx")
        doc = docx.Document(path)
        self._footnotes = self._load_footnotes(doc)

        page, img_n, tbl_n = 1, 0, 0
        buf: List[str] = []
        used_notes: List[str] = []

        def flush():
            nonlocal page, img_n, tbl_n, buf, used_notes
            if used_notes:
                buf.append("\n".join(used_notes))
            self.page_text[page] = _normalize_md("\n\n".join(b for b in buf if b))
            page, img_n, tbl_n, buf, used_notes = page + 1, 0, 0, [], []

        for block in self._iter_blocks(doc):
            if isinstance(block, Table):
                tbl_n += 1
                buf.append(f"## Table-{page}-{tbl_n}\n\n" + self._table_md(block))
                continue
            md, breaks, imgs, notes = self._paragraph_md(block, doc)
            for k in range(imgs):                      # image placeholders in-place
                img_n += 1
                stem = f"Figure-{page}-{img_n}"
                md = md.replace("\x00IMG\x00",
                                f"\n\n## {stem}\n\n", 1)
                self._save_para_image(block, stem, k + 1)
            used_notes += notes
            if md.strip():
                buf.append(md.strip())
            for _ in range(breaks):                    # explicit page breaks
                flush()
        flush()
        return self.page_text

    # -- internals ----------------------------------------------------------#

    @staticmethod
    def _iter_blocks(doc):
        """Yield paragraphs and tables in true document order."""
        from docx.document import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        body = doc.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, doc)
            elif child.tag == qn("w:tbl"):
                yield Table(child, doc)

    def _load_footnotes(self, doc) -> Dict[str, str]:
        notes: Dict[str, str] = {}
        try:
            part = doc.part.package.part_related_by(
                doc.part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes")
        except Exception:
            try:
                part = next(p for p in doc.part.package.iter_parts()
                            if p.partname.endswith("footnotes.xml"))
            except StopIteration:
                return notes
        from lxml import etree
        root = etree.fromstring(part.blob)
        for fn in root.findall("w:footnote", self._NS):
            fid = fn.get(f'{{{self._NS["w"]}}}id')
            text = " ".join(t.text or "" for t in fn.iter(f'{{{self._NS["w"]}}}t')).strip()
            if fid and text and int(fid) > 0:
                notes[fid] = text
        return notes

    def _paragraph_md(self, para, doc):
        """Return (markdown, page_break_count, image_count, footnote_lines)."""
        from docx.oxml.ns import qn
        style = (para.style.name or "").lower()

        parts: List[str] = []
        breaks = imgs = 0
        notes: List[str] = []

        def walk_run(r):
            nonlocal breaks, imgs
            out = []
            for node in r.iterchildren():
                tag = node.tag
                if tag == qn("w:t"):
                    out.append(node.text or "")
                elif tag == qn("w:tab"):
                    out.append("\t")
                elif tag == qn("w:br"):
                    if node.get(qn("w:type")) == "page":
                        breaks += 1
                    else:
                        out.append("\n")
                elif tag == qn("w:lastRenderedPageBreak"):
                    breaks += 1
                elif tag in (qn("w:drawing"), qn("w:pict"), qn("w:object")):
                    imgs += 1
                    out.append("\x00IMG\x00")
                elif tag == qn("w:footnoteReference"):
                    fid = node.get(qn("w:id"))
                    if fid in getattr(self, "_footnotes", {}):
                        out.append(f"[^{fid}]")
                        notes.append(f"[^{fid}]: {self._footnotes[fid]}")
            return self._fmt_run("".join(out), r)

        for child in para._p.iterchildren():
            if child.tag == qn("w:r"):
                parts.append(walk_run(child))
            elif child.tag == qn("w:hyperlink"):
                inner = "".join(walk_run(r) for r in child.findall(qn("w:r"))).strip()
                url = self._link_target(child, doc)
                parts.append(f"[{inner}]({url})" if url and inner else inner)

        text = "".join(parts)
        return self._decorate(text, para, style), breaks, imgs, notes

    @staticmethod
    def _fmt_run(text: str, r) -> str:
        """Apply bold/italic/inline-code markers from run properties."""
        if not text:
            return ""
        from docx.oxml.ns import qn
        rpr = r.find(qn("w:rPr"))
        if rpr is None or not text.strip():
            return text
        def on(tag):
            el = rpr.find(qn(tag))
            return el is not None and el.get(qn("w:val")) not in ("false", "0")
        body = text
        fonts = rpr.find(qn("w:rFonts"))
        mono = fonts is not None and any(
            "courier" in (fonts.get(qn(a)) or "").lower() or "consolas" in (fonts.get(qn(a)) or "").lower()
            for a in ("w:ascii", "w:hAnsi"))
        if mono:
            return f"`{body}`"
        if on("w:b") and body.strip():
            body = f"**{body.strip()}**" if body == body.strip() else body
        if on("w:i") and body.strip() and not body.startswith("**"):
            body = f"*{body.strip()}*"
        return body

    def _decorate(self, text: str, para, style: str) -> str:
        """Prefix heading/list/quote/code markers based on paragraph style."""
        from docx.oxml.ns import qn
        stripped = text.strip()
        if not stripped:
            return ""
        m = re.match(r"heading (\d+)", style)
        if m:
            return f"{'#' * min(int(m.group(1)) + 1, 6)} {stripped}"
        if style == "title":
            return f"# {stripped}"
        if "quote" in style:
            return f"> {stripped}"
        if any(k in style for k in ("code", "preformatted", "htmlpreformatted")):
            return f"```\n{stripped}\n```"
        if "caption" in style:
            return f"*{stripped}*"
        ppr = para._p.find(qn("w:pPr"))
        numpr = ppr.find(qn("w:numPr")) if ppr is not None else None
        if numpr is not None or "list" in style:
            ilvl = 0
            if numpr is not None:
                lvl = numpr.find(qn("w:ilvl"))
                ilvl = int(lvl.get(qn("w:val"))) if lvl is not None else 0
            marker = "1." if "number" in style else "-"
            if re.match(r"^\[( |x|X)\]", stripped):          # checklists
                marker = "-"
            return "  " * ilvl + f"{marker} {stripped}"
        return stripped

    def _link_target(self, node, doc) -> Optional[str]:
        from docx.oxml.ns import qn
        rid = node.get(qn("r:id"))
        if rid and rid in doc.part.rels:
            return doc.part.rels[rid].target_ref
        anchor = node.get(qn("w:anchor"))
        return f"#{anchor}" if anchor else None

    def _save_para_image(self, para, stem: str, seq: int):
        from docx.oxml.ns import qn
        blips = para._p.findall(".//" + qn("a:blip"))
        idx = min(seq, len(blips)) - 1
        if idx < 0 or idx >= len(blips):
            return
        rid = blips[idx].get(qn("r:embed")) or blips[idx].get(qn("r:link"))
        if rid and rid in para.part.rels:
            part = para.part.rels[rid].target_part
            ext = os.path.splitext(part.partname)[1]
            _save_image(part.blob, self.tags_dir, stem, ext)

    def _table_md(self, table) -> str:
        rows, seen = [], set()
        for row in table.rows:
            cells, prev_tc = [], None
            for cell in row.cells:
                if cell._tc is prev_tc:                      # merged cell repeat
                    continue
                prev_tc = cell._tc
                cells.append(cell.text)
            key = tuple(cells)
            if key in seen and not any(cells):               # drop duplicate blanks
                continue
            seen.add(key)
            rows.append(cells)
        return _md_table(rows)


# --------------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------------- #

class PPTExtractor:
    """Extract .ppt/.pptx to per-slide Markdown with figure/table tags."""

    def __init__(self, file_path: str, tags_dir: str):
        self.file_path = file_path
        self.tags_dir = tags_dir
        self.page_text: Dict[int, str] = {}

    def extract(self) -> Dict[int, str]:
        from pptx import Presentation

        path = self.file_path
        if path.lower().endswith(".ppt"):
            path = _convert_with_soffice(path, "pptx")
        prs = Presentation(path)

        for s_idx, slide in enumerate(prs.slides, start=1):
            state = {"img": 0, "tbl": 0}
            parts: List[str] = []
            title_id = slide.shapes.title.shape_id if slide.shapes.title else None

            for shape in self._reading_order(slide.shapes):
                md = self._shape_md(shape, s_idx, state, is_title=shape.shape_id == title_id)
                if md:
                    parts.append(md)

            notes = self._notes(slide)
            if notes:
                parts.append(f"> **Speaker notes:** {notes}")

            self.page_text[s_idx] = _normalize_md("\n\n".join(parts))
        return self.page_text

    # -- internals ----------------------------------------------------------#

    @staticmethod
    def _reading_order(shapes):
        """Top-to-bottom, left-to-right ordering for natural reading flow."""
        def key(sh):
            top = sh.top if sh.top is not None else 0
            left = sh.left if sh.left is not None else 0
            return (top, left)
        return sorted(shapes, key=key)

    def _shape_md(self, shape, s_idx: int, state: dict, is_title=False) -> str:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            return "\n\n".join(
                filter(None, (self._shape_md(sub, s_idx, state)
                              for sub in self._reading_order(shape.shapes))))
        if st == MSO_SHAPE_TYPE.PICTURE or getattr(shape, "image", None) is not None:
            return self._figure(shape, s_idx, state)
        if shape.has_table:
            state["tbl"] += 1
            rows = [[c.text for c in row.cells] for row in shape.table.rows]
            return f"## Table-{s_idx}-{state['tbl']}\n\n" + _md_table(rows)
        if getattr(shape, "has_chart", False) and shape.has_chart:
            return self._chart(shape, s_idx, state)
        if st == MSO_SHAPE_TYPE.LINKED_PICTURE or st == MSO_SHAPE_TYPE.MEDIA:
            state["img"] += 1
            return f"## Figure-{s_idx}-{state['img']}"
        if shape.has_text_frame:
            return self._text_frame_md(shape.text_frame, is_title)
        return ""

    def _figure(self, shape, s_idx: int, state: dict) -> str:
        state["img"] += 1
        stem = f"Figure-{s_idx}-{state['img']}"
        try:
            img = shape.image
            _save_image(img.blob, self.tags_dir, stem, img.ext)
        except Exception:
            pass
        return f"## {stem}"

    def _chart(self, shape, s_idx: int, state: dict) -> str:
        """Charts: emit figure tag + underlying data table for LLM usability."""
        state["img"] += 1
        stem = f"Figure-{s_idx}-{state['img']}"
        md = [f"## {stem}"]
        try:
            chart = shape.chart
            cats = [str(c) for c in chart.plots[0].categories]
            rows = [["Category"] + [s.name or f"Series {i+1}"
                                    for i, s in enumerate(chart.plots[0].series)]]
            series_vals = [list(s.values) for s in chart.plots[0].series]
            for i, cat in enumerate(cats):
                rows.append([cat] + [("" if i >= len(v) or v[i] is None else str(v[i]))
                                     for v in series_vals])
            title = chart.chart_title.text_frame.text if chart.has_title else ""
            if title:
                md.append(f"*Chart: {title}*")
            md.append(_md_table(rows))
        except Exception:
            pass
        return "\n\n".join(md)

    @staticmethod
    def _text_frame_md(tf, is_title: bool) -> str:
        lines: List[str] = []
        for p in tf.paragraphs:
            text = "".join(PPTExtractor._run_md(r) for r in p.runs).strip()
            if not text:
                continue
            if is_title and not lines:
                lines.append(f"# {text}")
            elif p.level > 0 or (len(tf.paragraphs) > 1 and not is_title):
                lines.append("  " * max(p.level, 0) + f"- {text}")
            else:
                lines.append(text)
        return "\n".join(lines)

    @staticmethod
    def _run_md(run) -> str:
        text = run.text
        if not text.strip():
            return text
        try:
            if run.font.bold:
                text = f"**{text.strip()}**"
            if run.font.italic and not text.startswith("**"):
                text = f"*{text.strip()}*"
            link = run.hyperlink.address
            if link:
                text = f"[{text}]({link})"
        except Exception:
            pass
        return text

    @staticmethod
    def _notes(slide) -> str:
        if not slide.has_notes_slide:
            return ""
        tf = slide.notes_slide.notes_text_frame
        return _WS_RE.sub(" ", tf.text.strip()) if tf is not None else ""


# --------------------------------------------------------------------------- #
# Excel / CSV
# --------------------------------------------------------------------------- #

class ExcelExtractor:
    """Extract .csv/.xls/.xlsx to per-sheet Markdown with figure/table tags."""

    def __init__(self, file_path: str, tags_dir: str):
        self.file_path = file_path
        self.tags_dir = tags_dir
        self.page_text: Dict[int, str] = {}

    def extract(self) -> Dict[int, str]:
        ext = os.path.splitext(self.file_path)[1].lower()
        if ext == ".csv":
            return self._extract_csv()
        path = self.file_path
        if ext == ".xls":
            path = _convert_with_soffice(path, "xlsx")
        return self._extract_xlsx(path)

    # -- CSV ---------------------------------------------------------------- #

    def _extract_csv(self) -> Dict[int, str]:
        name = os.path.splitext(os.path.basename(self.file_path))[0]
        with open(self.file_path, newline="", encoding="utf-8-sig", errors="replace") as fh:
            sample = fh.read(8192); fh.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
            except csv.Error:
                dialect = csv.excel
            rows = [[c for c in row] for row in csv.reader(fh, dialect)]
        rows = [r for r in rows if any(str(c).strip() for c in r)]
        md = f"# Sheet: {name}\n\n## Table-{name}-1\n\n{_md_table(rows)}"
        self.page_text[1] = _normalize_md(md)
        return self.page_text

    # -- XLSX --------------------------------------------------------------- #

    def _extract_xlsx(self, path: str) -> Dict[int, str]:
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)

        for s_idx, ws in enumerate(wb.worksheets, start=1):
            name = ws.title
            parts = [f"# Sheet: {name}"]
            img_n = 0
            for img in getattr(ws, "_images", []):
                img_n += 1
                stem = f"Figure-{name}-{img_n}"
                try:
                    blob = img._data() if callable(getattr(img, "_data", None)) else img.ref.getvalue()
                    fmt = getattr(img, "format", None) or "png"
                    _save_image(blob, self.tags_dir, stem, fmt)
                except Exception:
                    pass
                parts.append(f"## {stem}")
            for _ in getattr(ws, "_charts", []):
                img_n += 1
                parts.append(f"## Figure-{name}-{img_n}")

            for t_idx, block in enumerate(self._table_blocks(ws), start=1):
                parts.append(f"## Table-{name}-{t_idx}\n\n" + _md_table(block))

            self.page_text[s_idx] = _normalize_md("\n\n".join(parts))
        return self.page_text

    @staticmethod
    def _cell_str(v) -> str:
        if v is None:
            return ""
        if isinstance(v, float) and v == int(v):
            return str(int(v))
        return str(v)

    def _table_blocks(self, ws) -> List[List[List[str]]]:
        """Split the used range into contiguous row blocks (blank rows = separators)."""
        blocks, current = [], []
        for row in ws.iter_rows(values_only=True):
            vals = [self._cell_str(v) for v in row]
            if any(v.strip() for v in vals):
                current.append(vals)
            elif current:
                blocks.append(self._trim(current))
                current = []
        if current:
            blocks.append(self._trim(current))
        return blocks

    @staticmethod
    def _trim(block: List[List[str]]) -> List[List[str]]:
        """Drop fully-empty trailing columns from a block."""
        width = 0
        for row in block:
            for i in range(len(row) - 1, -1, -1):
                if row[i].strip():
                    width = max(width, i + 1)
                    break
        return [row[:width] for row in block]


# --------------------------------------------------------------------------- #
# Dispatch convenience
# --------------------------------------------------------------------------- #

_EXT_MAP = {
    ".doc": WordExtractor, ".docx": WordExtractor,
    ".ppt": PPTExtractor, ".pptx": PPTExtractor,
    ".csv": ExcelExtractor, ".xls": ExcelExtractor, ".xlsx": ExcelExtractor,
}


def extract(file_path: str, tags_dir: str) -> Dict[int, str]:
    """Route any supported Office file to the right extractor."""
    ext = os.path.splitext(file_path)[1].lower()
    cls = _EXT_MAP.get(ext)
    if cls is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return cls(file_path, tags_dir).extract()


if __name__ == "__main__":
    import json
    import sys
    pages = extract(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "./tags")
    for n, content in pages.items():
        print(f"\n{'=' * 20} PAGE/SLIDE/SHEET {n} {'=' * 20}\n{content}")
